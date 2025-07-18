import os
import json
from datetime import datetime
from aiohttp import web
from telegram import Update
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes
from openai import OpenAI
import asyncio
import gspread
from oauth2client.service_account import ServiceAccountCredentials
import aiohttp
import random
# مكتبات قراءة الملفات
import fitz  # PyMuPDF
from docx import Document as DocxReader
from pptx import Presentation
import pandas as pd
import tempfile

# ========== المتغيرات ==========
BOT_TOKEN = os.getenv("BOT_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_CX = os.getenv("GOOGLE_CX")
IMGBB_API_KEY = os.getenv("IMGBB_API_KEY") 
MAX_SESSION_LENGTH = 20

client = OpenAI(api_key=OPENAI_API_KEY)
group_sessions = {}
group_dialects = {}
image_context = {}  # تخزين مؤقت للصورة لكل مستخدم

# ========== Google Sheets ==========
scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
creds = ServiceAccountCredentials.from_json_keyfile_name("service_account.json", scope)
gc = gspread.authorize(creds)
sheet = gc.open("RahimBot_History").sheet1

# ========== تحميل برومبت ==========
with open("rahim_prompts_library.txt", "r", encoding="utf-8") as f:
    PROMPTS_LIBRARY = f.read()

RAHIM_MAIN_PROMPT = (
    "أنت بوت اسمه \"رحيم\". تم تصميمك لتكون زول طيب، حنون، خفيف الدم، وبتتكلم بلغة بشرية حقيقية. \
    لو حسيت إنو السائل محتاج معلومة دقيقة، ممكن تفتش في الإنترنت وترد عليه برابط موثوق."
)

# ========== أدوات قراءة الملفات ==========
def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_docx(file_path):
    doc = DocxReader(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        return f"📛 حصل خطأ أثناء قراءة ملف PowerPoint: {str(e)}"

def extract_text_from_excel(file_path):
    try:
        df = pd.read_excel(file_path)
        preview = df.head(10)
        return preview.to_string(index=False)
    except Exception as e:
        return f"📛 حصل خطأ أثناء قراءة الإكسل: {str(e)}"

# ========== حفظ في Google Sheets ==========
def save_message_to_sheet(data):
    try:
        sheet.append_row([
            data["timestamp"],
            str(data["user_id"]),
            data["user_name"],
            str(data["group_id"]),
            data["dialect"],
            data["text"]
        ])
        print("✅ Saved to Google Sheet", flush=True)
    except Exception as e:
        print(f"❌ Error saving to Google Sheet: {e}", flush=True)

# ========== كشف اللهجة ==========
async def detect_language_or_dialect(text: str) -> str:
    prompt = (
        "حدد لي لغة أو لهجة النص التالي بدقة عالية..."
        f"النص: \"{text}\""
    )
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "أنت مساعد لتحديد لهجة أو لغة النص."},
                {"role": "user", "content": prompt}
            ]
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"Error detecting dialect/language: {e}", flush=True)
        return "العربية الفصحى"

# ========== البحث في الويب ==========
async def perform_web_search(query: str) -> str:
    try:
        async with aiohttp.ClientSession() as session:
            url = f"https://www.googleapis.com/customsearch/v1?key={GOOGLE_API_KEY}&cx={GOOGLE_CX}&q={query}"
            async with session.get(url) as response:
                data = await response.json()
                if "items" in data and len(data["items"]) > 0:
                    top = data["items"][0]
                    title = top["title"]
                    snippet = top["snippet"]
                    link = top["link"]
                    return f"أنا مش مختص، لكن لقيت ليك من Google:\n**{title}**\n{snippet}\n📎 {link}"
                else:
                    return "ما لقيت نتيجة واضحة في البحث 😕"
    except Exception as e:
        return f"📛 حصل خطأ أثناء البحث: {str(e)}"

# ========== دوال ترفيه ومحادثة عامة ==========
async def call_openai_chat(prompt: str, max_tokens=500, temperature=0.7) -> str:
    try:
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=max_tokens,
            temperature=temperature,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"OpenAI error in call_openai_chat: {e}")
        return "حصل خطأ أثناء التواصل مع الذكاء الصناعي."

# أمر /أغنية
async def suggest_song(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = " ".join(context.args).strip()
    if not query:
        await update.message.reply_text("اكتب لي نوع أو مزاج الأغنية عشان أساعدك.")
        return
    prompt = f"اقترح لي قائمة أغاني تناسب المزاج أو النوع التالي: {query}"
    response = await call_openai_chat(prompt)
    await update.message.reply_text(response)

# أمر /كتاب
async def suggest_book(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = " ".join(context.args).strip()
    if not query:
        await update.message.reply_text("اكتب لي نوع الكتاب اللي تحبه.")
        return
    prompt = f"اقترح لي كتب مميزة من نوع: {query} مع ملخص بسيط لكل كتاب."
    response = await call_openai_chat(prompt)
    await update.message.reply_text(response)

# أمر /فيلم
async def suggest_movie(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = " ".join(context.args).strip()
    if not query:
        await update.message.reply_text("اكتب لي نوع الفيلم اللي تحبه.")
        return
    prompt = f"اقترح لي أفلام مميزة من نوع: {query} مع ملخص بسيط لكل فيلم."
    response = await call_openai_chat(prompt)
    await update.message.reply_text(response)

# أمر /نقاش
async def start_discussion(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = " ".join(context.args).strip()
    if not topic:
        await update.message.reply_text("اكتب لي موضوع للنقاش.")
        return
    prompt = f"ابدأ نقاش ترفيهي عن الموضوع التالي: {topic}. اسأل المتابعين أسئلة تحفزهم على التفاعل."
    response = await call_openai_chat(prompt)
    await update.message.reply_text(response)

# أمر /لعبة
brain_games = [
    "فكر في رقم بين 1 و 10",
    "حل اللغز: ما هو الشيء الذي له أسنان ولا يعض؟",
    "إذا كان اليوم الإثنين، ما هو اليوم بعد 100 يوم؟",
]

async def play_game(update: Update, context: ContextTypes.DEFAULT_TYPE):
    game = random.choice(brain_games)
    await update.message.reply_text(f"لعبتنا اليوم: {game}")

# أمر /سؤال
async def answer_question(update: Update, context: ContextTypes.DEFAULT_TYPE):
    question = " ".join(context.args).strip()
    if not question:
        await update.message.reply_text("اكتب لي السؤال اللي عايز تعرف إجابته.")
        return
    prompt = f"جاوب على السؤال التالي ببساطة ووضوح: {question}"
    response = await call_openai_chat(prompt)
    await update.message.reply_text(response)

# ردود ودية عشوائية على الرسائل العادية (غير أوامر)
async def friendly_reply(update: Update, context: ContextTypes.DEFAULT_TYPE):
    texts = [
        "جميل الكلام دا 😊",
        "حلو شديد! 😊",
        "مبسوط من تواصلك معانا! 🌟",
        "دا كلام جميل، خلي نواصل 😊",
    ]
    if update.message.text and not update.message.text.startswith("/"):
        await update.message.reply_text(random.choice(texts))
        
# ========== /start ==========
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    group_id = update.message.chat.id
    dialect = "العربية الفصحى"
    group_dialects[group_id] = dialect
    combined_prompt = RAHIM_MAIN_PROMPT + "\n\n" + PROMPTS_LIBRARY
    group_sessions[group_id] = [{"role": "system", "content": combined_prompt}]
    await update.message.reply_text("البوت شغال ✅")

from trends_manager import get_next_trend_lifo  # استدعاء دالة الترند الجديدة

# ========== رسائل القروب ==========
async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    bot_username = (await context.bot.get_me()).username.lower()
    user_message = update.message.text.lower()

    if (
        f"@{bot_username}" not in user_message
        and "رحيم" not in user_message
        and "rahim" not in user_message
        and not update.message.reply_to_message
    ):
        return

    group_id = update.message.chat.id
    user_id = update.message.from_user.id
    user_name = update.message.from_user.full_name

    # لو المستخدم طلب ترند صريح
    if user_message in ["/trend", "ادينا ترند"]:
        group_type = get_group_type(group_id)  # لازم تعملي الدالة دي حسب القروب
        trend = get_next_trend_lifo(group_type)
        if trend:
            await context.bot.send_message(chat_id=group_id, text=trend)
        return  # بعد ما أرسل الترند نوقف

    # استكمال باقي الكود العادي
    if update.message.reply_to_message and update.message.reply_to_message.text:
        target_text = update.message.reply_to_message.text
        combined_input = f"{update.message.text}\n\nالرسالة المردود عليها:\n{target_text}"
    else:
        combined_input = update.message.text

    if group_id not in group_sessions:
        detected = await detect_language_or_dialect(combined_input)
        group_dialects[group_id] = detected
        combined_prompt = RAHIM_MAIN_PROMPT + "\n\n" + PROMPTS_LIBRARY
        group_sessions[group_id] = [{"role": "system", "content": combined_prompt}]

    else:
        detected = group_dialects.get(group_id, "العربية الفصحى")

    group_sessions[group_id].append({"role": "user", "content": combined_input})
    group_sessions[group_id] = group_sessions[group_id][-MAX_SESSION_LENGTH:]

    save_message_to_sheet({
        "user_id": user_id,
        "user_name": user_name,
        "group_id": group_id,
        "text": combined_input,
        "dialect": detected,
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    })

    # ======= لو في كلمات طبية، يتم البحث أولاً =======
    if any(x in combined_input for x in ["علاج", "تشخيص", "أعراض", "مرض", "دواء"]):
        web_result = await perform_web_search(combined_input)
        if "ما لقيت نتيجة واضحة" not in web_result and "📛 حصل خطأ" not in web_result:
            await update.message.reply_text(web_result)
            return

    try:
        model = "gpt-4o" if "معلومة دقيقة" in combined_input else "gpt-3.5-turbo"
        response = client.chat.completions.create(
            model=model,
            messages=group_sessions[group_id],
            max_tokens=2000,
        )
        full_reply = response.choices[0].message.content.strip()
        group_sessions[group_id].append({"role": "assistant", "content": full_reply})
        group_sessions[group_id] = group_sessions[group_id][-MAX_SESSION_LENGTH:]

        # تقسيم الرد الطويل إلى دفعات
        async def send_in_chunks(text, chunk_size=1500):
            chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
            for chunk in chunks:
                await update.message.reply_text(chunk)
                await asyncio.sleep(1)

        await send_in_chunks(full_reply)

    except Exception as e:
        await update.message.reply_text("حصل خطأ في الذكاء الصناعي 😔")
        print(f"OpenAI error: {e}", flush=True)

# ========== الخاص ==========
async def handle_private_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(
        "مرحباً! 👋 البوت دا مخصص للقروبات فقط.\n"
        "Please note: This bot is designed for group chats only.\n"
        "أضفني لقروبك عشان أقدر أساعدك 🚀"
    )

# ========== التعامل مع الملفات ==========
async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    file = await update.message.document.get_file()
    file_ext = file.file_path.split('.')[-1].lower()

    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_ext}") as temp:
        await file.download_to_drive(temp.name)

        if file_ext == "pdf":
            text = extract_text_from_pdf(temp.name)
        elif file_ext in ["docx", "doc"]:
            text = extract_text_from_docx(temp.name)
        elif file_ext in ["xlsx", "xls"]:
            text = extract_text_from_excel(temp.name)
        elif file_ext == "pptx":
            text = extract_text_from_pptx(temp.name)
        else:
            await update.message.reply_text("حالياً بقدر أقرأ ملفات PDF, Word, Excel و PowerPoint فقط.")
            return

        if len(text) > 2000:
            text = text[:2000] + "\n\n📌 النص طويل جداً، تم عرض جزء فقط."

        await update.message.reply_text(
            f"📂 تم استخراج المحتوى:\n\n{text}\n\nتحب أعمل شنو بيهو؟ (تلخيص؟ تحليل؟ فلاش كارد؟)"
        )

# ========== Webhook ==========
async def webhook(request):
    try:
        data = await request.json()
        update = Update.de_json(data, application.bot)
        await application.process_update(update)
    except Exception as e:
        print(f"Webhook error: {e}", flush=True)
    return web.Response(text="OK")

# ========== تشغيل التطبيق ==========
# ========== تحديث إضافة المعالجات للأوامر في main ==========
application = Application.builder().token(BOT_TOKEN).build()

# الأوامر الأساسية الموجودة
application.add_handler(CommandHandler("start", start))
application.add_handler(MessageHandler(filters.TEXT & filters.ChatType.GROUPS & ~filters.COMMAND, handle_message))
application.add_handler(MessageHandler(filters.TEXT & filters.ChatType.PRIVATE & ~filters.COMMAND, handle_private_message))
application.add_handler(MessageHandler(filters.Document.ALL, handle_document))

# إضافة أوامر الترفيه والمحادثة
application.add_handler(CommandHandler("song", suggest_song))
application.add_handler(CommandHandler("book", suggest_book))
application.add_handler(CommandHandler("movie", suggest_movie))
application.add_handler(CommandHandler("discussion", start_discussion))
application.add_handler(CommandHandler("game", play_game))
application.add_handler(CommandHandler("question", answer_question))

application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, friendly_reply))
app = web.Application()
app.router.add_post(f'/{BOT_TOKEN}', webhook)

async def run():
    await application.initialize()
    await application.start()
    runner = web.AppRunner(app)
    await runner.setup()
    site = web.TCPSite(runner, port=int(os.getenv("PORT", 8080)))
    await site.start()
    print("💬 البوت شغال على السيرفر...")
    await asyncio.Event().wait()

asyncio.run(run())
